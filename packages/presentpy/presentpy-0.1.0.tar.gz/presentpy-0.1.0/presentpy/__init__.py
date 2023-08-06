import sys
from pathlib import Path

import nbformat
import mistletoe
from pptx import Presentation

import pygments
import pygments.styles
from pygments import lex
from pygments.lexers.python import PythonLexer
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pygments.token import Token

documents = []

presentation = Presentation()

style = pygments.styles.get_style_by_name("default")
token_colors = {}
for token, str_style in style.styles.items():
    if not str_style:
        continue
    _, _, color = str_style.partition("#")
    if not color:
        continue

    pad = 1 if len(color) == 3 else 2
    token_colors[token] = RGBColor(*[int(color[i : i + pad], 16) for i in range(0, len(color), pad)])
token_colors[Token.Keyword] = RGBColor(0, 128, 20)
token_colors[Token.Name.Class] = RGBColor(49, 0, 250)
token_colors[Token.Name.Builtin.Pseudo] = RGBColor(27, 82, 167)
token_colors[Token.Keyword.Constant] = token_colors[Token.Keyword]
token_colors[Token.Name.Function.Magic] = token_colors[Token.Name.Class]
token_colors[Token.Name] = RGBColor(27, 82, 167)
token_colors[Token.Comment.Single] = RGBColor(76, 135, 135)
token_colors[Token.Operator] = RGBColor(175, 24, 251)


def add_bullet_slide(prs, title, bullet_points):
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    title_shape.text = title

    body_shape = shapes.placeholders[1]

    tf = body_shape.text_frame

    for bullet_point in bullet_points:
        p = tf.add_paragraph()
        p.text = bullet_point
        p.level = 1


def add_code_slide(prs, title, code):
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    title_shape.text = title

    body_shape = shapes.placeholders[1]

    text_frame = body_shape.text_frame

    text_frame.clear()

    p = text_frame.paragraphs[0]
    p.bullet = False
    for kind, content in code:
        run = p.add_run()
        run.bullet = False
        if content == "\n":
            run.text = "\x0A"
        else:
            run.text = content

        font = run.font
        font.color.rgb = token_colors.get(kind, RGBColor(0, 0, 0))
        font.name = "Courier"
        font.size = Pt(18)


file = Path(sys.argv[1])

with open(file) as r:
    notebook = nbformat.read(r, as_version=4)

    for cell in notebook["cells"]:
        if cell["cell_type"] == "markdown":
            source = cell["source"]
            document = mistletoe.Document(source)
            header = document.children[0]
            bullets = []
            if len(document.children) > 1:
                bullets = [bullet.children[0].children[0].content for bullet in document.children[1].children]
            add_bullet_slide(
                presentation,
                header.children[0].content,
                bullets,
            )

        elif cell["cell_type"] == "code":
            source = cell["source"]
            if not source:
                continue

            rtf = lex(source, PythonLexer())

            add_code_slide(
                presentation,
                "code",
                rtf,
            )

presentation.save(f"{file.stem}.pptx")

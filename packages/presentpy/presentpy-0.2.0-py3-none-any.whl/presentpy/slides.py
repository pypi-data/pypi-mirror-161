from typing import Any, Iterable, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.util import Pt

from presentpy.code import get_styles
from presentpy.code_cell_config import CodeCellConfig

token_colors = get_styles()


def add_title_slide(prs: Presentation, title: str, subtitle: Optional[str] = None):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    if subtitle:
        subtitle_shape.text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullet_points: List[str]) -> None:
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    title_shape.text = title

    body_shape = shapes.placeholders[1]

    tf = body_shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    tf.text = bullet_points[0]

    for bullet_point in bullet_points[1:]:
        p = tf.add_paragraph()
        p.line_spacing = 1.5
        p.text = bullet_point
        p.level = 0


def add_code_slide(prs: Presentation, parsed_lines: List[List[Tuple[Any, str]]], config: CodeCellConfig) -> None:
    highlights = config.highlights
    if not highlights:
        highlights = [[0]]
    else:
        highlights = [[0]] + highlights
    for hl in highlights:
        add_code_slide_highlighted(prs, parsed_lines, config.title, highlights=hl)


def add_code_slide_highlighted(
    prs: Presentation, parsed_lines: List[List[Tuple[Any, str]]], title: Optional[str], highlights: Iterable[int]
) -> None:
    highlighted_lines = set(highlights)

    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    if title:
        title_shape = shapes.title
        title_shape.text = title

    body_shape = shapes.placeholders[1]

    text_frame = body_shape.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    p = text_frame.paragraphs[0]
    p.bullet = False

    for ln, line in enumerate(parsed_lines, 1):
        for kind, text in line:
            run = p.add_run()
            run.text = text
            font = run.font
            font.bold = ln in highlighted_lines
            font.color.rgb = token_colors.get(kind, RGBColor(0, 0, 0))
            font.name = "Courier"
            font.size = Pt(14)
        run = p.add_run()
        run.font.size = Pt(14)
        run.text = "\x0A"
